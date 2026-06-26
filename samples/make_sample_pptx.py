#!/usr/bin/env python3
"""
make_sample_pptx.py — Generate a representative academic-medical .pptx so the
extract -> build -> compare pipeline can run without a proprietary deck.

Produces figs/ placeholder images (CT-like, 4-panel, logos) and sample.pptx
with: title slide, single-figure slide, grouped multi-panel slide, table slide,
each carrying institutional chrome (logos, footer, bottom blue line).

Requires: python-pptx, pillow
"""
import os
from PIL import Image, ImageDraw

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
FIGS = os.path.join(HERE, "figs")
os.makedirs(FIGS, exist_ok=True)

BRAND = RGBColor(0x1F, 0x4E, 0x79)   # institutional blue


# --------------------------- placeholder figures ---------------------------
def make_ct(path, w=1400, h=1100):
    img = Image.new("RGB", (w, h), (8, 8, 8))
    d = ImageDraw.Draw(img)
    cx, cy = w // 2, h // 2
    for r in range(min(cx, cy), 0, -2):
        v = int(220 * (r / min(cx, cy)) ** 0.6)
        d.ellipse([cx - r, cy - r, cx + r, cy + r], outline=(v, v, v))
    d.ellipse([cx - 90, cy - 60, cx - 10, cy + 40], fill=(40, 40, 40))
    d.ellipse([cx + 20, cy - 50, cx + 110, cy + 50], fill=(60, 60, 60))
    img.save(path)


def make_panel(path, tint, w=700, h=520):
    img = Image.new("RGB", (w, h), tint)
    d = ImageDraw.Draw(img)
    for i in range(0, w, 14):
        d.line([(i, 0), (i, h)], fill=tuple(min(255, c + 18) for c in tint))
    d.ellipse([w//2-120, h//2-90, w//2+120, h//2+90], outline=(255, 255, 255), width=3)
    img.save(path)


def make_logo(path, text, color):
    w, h = 360, 150
    img = Image.new("RGBA", (w, h), (255, 255, 255, 0))
    d = ImageDraw.Draw(img)
    d.rounded_rectangle([4, 4, w-4, h-4], radius=16, fill=color)
    d.text((24, 54), text, fill=(255, 255, 255))
    img.save(path)


make_ct(os.path.join(FIGS, "ct.png"))
for nm, tint in [("p_a", (60, 30, 30)), ("p_b", (30, 50, 40)),
                 ("p_c", (35, 35, 60)), ("p_d", (60, 55, 30))]:
    make_panel(os.path.join(FIGS, nm + ".png"), tint)
make_logo(os.path.join(FIGS, "logo_hosp.png"), "YU HOSPITAL", (0x1F, 0x4E, 0x79))
make_logo(os.path.join(FIGS, "logo_univ.png"), "YEUNGNAM UNIV", (0x7A, 0x12, 0x12))


# --------------------------- deck construction ---------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

TITLE_ONLY = prs.slide_layouts[5]
BLANK = prs.slide_layouts[6]


def add_chrome(slide):
    """Institutional chrome present on content slides: logos, footer, blue line."""
    slide.shapes.add_picture(os.path.join(FIGS, "logo_hosp.png"),
                             Inches(0.3), Inches(0.25), Inches(1.5), Inches(0.62))
    slide.shapes.add_picture(os.path.join(FIGS, "logo_univ.png"),
                             Inches(11.5), Inches(0.25), Inches(1.5), Inches(0.62))
    # bottom blue line (thin rectangle)
    line = slide.shapes.add_shape(1, Inches(0), Inches(6.95), Inches(13.333), Pt(4))
    line.fill.solid(); line.fill.fore_color.rgb = BRAND
    line.line.fill.background()
    # footer
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "2026 대한심장학회 춘계학술대회  |  영남대학교병원 순환기내과"
    r.font.size = Pt(11); r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def add_title(slide, text):
    slide.shapes.title.text = text
    tf = slide.shapes.title.text_frame.paragraphs[0]
    tf.font.size = Pt(32); tf.font.bold = True
    tf.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)


def add_citation(slide, text):
    tb = slide.shapes.add_textbox(Inches(8.3), Inches(6.55), Inches(4.8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = text
    r.font.size = Pt(12); r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


# Slide 1 — Title
s = prs.slides.add_slide(prs.slide_layouts[0])
s.shapes.title.text = "관상동맥 석회화와 심혈관 예후"
s.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
s.placeholders[1].text = "전향적 코호트 연구 · Single-center experience"
add_chrome(s)

# Slide 2 — Single figure + caption + citation
s = prs.slides.add_slide(TITLE_ONLY)
add_title(s, "관상동맥 CT: 좌전하행지 석회화")
s.shapes.add_picture(os.path.join(FIGS, "ct.png"), Inches(3.2), Inches(1.3),
                     Inches(6.9), Inches(4.6))
cap = s.shapes.add_textbox(Inches(3.2), Inches(5.95), Inches(6.9), Inches(0.5))
cp = cap.text_frame.paragraphs[0]
r1 = cp.add_run(); r1.text = "Figure 1. Axial CT showing "
r1.font.size = Pt(16); r1.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
r2 = cp.add_run(); r2.text = "Agatston"; r2.font.italic = True
r2.font.size = Pt(16); r2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
r3 = cp.add_run(); r3.text = " score 312."
r3.font.size = Pt(16); r3.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
add_citation(s, "Circulation. 2023;147:1234-45.")
add_chrome(s)

# Slide 3 — Grouped 2x2 multi-panel figure
s = prs.slides.add_slide(TITLE_ONLY)
add_title(s, "다중 영상 비교 (A–D)")
panels = [("p_a", "A"), ("p_b", "B"), ("p_c", "C"), ("p_d", "D")]
positions = [(3.6, 1.4), (7.4, 1.4), (3.6, 3.9), (7.4, 3.9)]
shapes_to_group = []
for (nm, lab), (x, y) in zip(panels, positions):
    pic = s.shapes.add_picture(os.path.join(FIGS, nm + ".png"),
                               Inches(x), Inches(y), Inches(3.5), Inches(2.3))
    shapes_to_group.append(pic)
    lb = s.shapes.add_textbox(Inches(x + 0.05), Inches(y + 0.02), Inches(0.5), Inches(0.4))
    lr = lb.text_frame.paragraphs[0].add_run(); lr.text = lab
    lr.font.bold = True; lr.font.size = Pt(20)
    shapes_to_group.append(lb)
add_citation(s, "Eur Heart J. 2024;45:880-9.")
add_chrome(s)

# Slide 4 — Clinical table
s = prs.slides.add_slide(TITLE_ONLY)
add_title(s, "기저 특성 (Baseline characteristics)")
rows, cols = 4, 3
tbl = s.shapes.add_table(rows, cols, Inches(2.2), Inches(1.6),
                         Inches(8.9), Inches(3.2)).table
hdr = ["변수", "CAC < 100 (n=210)", "CAC ≥ 100 (n=188)"]
data = [["연령 (years)", "58.2 ± 9.1", "66.4 ± 8.3"],
        ["당뇨 (%)", "18.6", "34.0"],
        ["LDL-C (mg/dL)", "112 ± 31", "118 ± 29"]]
for c, h in enumerate(hdr):
    tbl.cell(0, c).text = h
for ri, row in enumerate(data, start=1):
    for ci, val in enumerate(row):
        tbl.cell(ri, ci).text = val
add_citation(s, "J Am Coll Cardiol. 2023;81:55-66.")
add_chrome(s)

out = os.path.join(HERE, "sample.pptx")
prs.save(out)
print("Wrote", out, "with", len(prs.slides._sldIdLst), "slides")
