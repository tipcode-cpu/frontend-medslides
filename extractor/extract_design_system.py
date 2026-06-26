#!/usr/bin/env python3
"""
extract_design_system.py — design system + chrome + objects + images + media,
with FULL text-style resolution (per-paragraph/per-run effective properties).

Text fidelity: PowerPoint body text inherits size, bullet, indentation and
spacing from the slide MASTER's bodyStyle/titleStyle by paragraph LEVEL. A flat
".ds-text" loses all of that. This extractor resolves, for every paragraph and
run, the EFFECTIVE: font size, bold, color, font family, alignment, bullet
char+font, left margin/indent, line spacing and paragraph spacing — by merging
run rPr > paragraph pPr > placeholder lstStyle > master style[level].

Also: master/layout chrome (logo, blue line, footer), robust image pipeline
(originals+display, raster EMF/WMF/WDP, never skip), media engine (video +
poster, never silent), and assets/extraction-report.json.

Usage: python extract_design_system.py <input.pptx> <output_dir>
"""
import io
import json
import os
import shutil
import sys
from collections import Counter

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from PIL import Image
    _HAS_PIL = True
except Exception:
    _HAS_PIL = False

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
NS = {"a": A[1:-1], "p": P[1:-1], "r": R[1:-1]}

EMU_PT = 12700
RASTER_OK = {"png", "jpg", "jpeg", "gif", "webp"}
NEED_RASTER = {"wmf", "emf", "wdp", "tiff", "tif", "bmp"}
DISPLAY_MAXDIM = 2600
DISPLAY_BIG_BYTES = 6_000_000
BROWSER_VIDEO = {"mp4", "m4v", "webm", "ogg", "ogv"}
# Symbol-font bullet codes -> Unicode (portable; avoids needing Wingdings)
WINGDINGS = {0x76: "❖", 0xa7: "▪", 0x2022: "•", 0x2013: "–",
             0xbb: "»", 0xa8: "▪", 0xfc: "✔", 0x9f: "•"}
ALGN = {"l": "left", "ctr": "center", "r": "right", "just": "justify"}

REPORT = {"images": [], "videos": [], "logos": [], "failures": []}
MASTER_LEVELS = {"title": {}, "body": {}}      # set in main()


# ----------------------------- color / geometry -----------------------------
def rgb_or_none(color):
    try:
        return ("#" + str(color.rgb).lower()) if color is not None else None
    except Exception:
        return None


def frac(l, t, w, h, W, H):
    return {"xf": (l or 0) / W, "yf": (t or 0) / H, "wf": (w or 0) / W, "hf": (h or 0) / H}


def shape_frac(sh, W, H):
    try:
        return frac(int(sh.left or 0), int(sh.top or 0), int(sh.width or 0), int(sh.height or 0), W, H)
    except Exception:
        return {"xf": 0, "yf": 0, "wf": 0, "hf": 0}


def group_child_transform(group):
    try:
        x = group._element.find(".//a:xfrm", NS)
        off, ext, choff, chext = x.find("a:off", NS), x.find("a:ext", NS), x.find("a:chOff", NS), x.find("a:chExt", NS)
        ox, oy, ex, ey = int(off.get("x")), int(off.get("y")), int(ext.get("cx")), int(ext.get("cy"))
        cox, coy, cex, cey = int(choff.get("x")), int(choff.get("y")), int(chext.get("cx")), int(chext.get("cy"))
        sx, sy = (ex / cex if cex else 1.0), (ey / cey if cey else 1.0)
        return lambda l, t, w, h: (ox + (l - cox) * sx, oy + (t - coy) * sy, w * sx, h * sy)
    except Exception:
        return lambda l, t, w, h: (l, t, w, h)


# ----------------------------- text style resolution -----------------------------
def parse_props(el):
    """Effective paragraph-level props from a pPr / lvlNpPr element."""
    if el is None:
        return {}
    p = {}
    if el.get("marL") is not None:
        p["marL"] = int(el.get("marL")) / EMU_PT
    if el.get("indent") is not None:
        p["indent"] = int(el.get("indent")) / EMU_PT
    if el.get("algn"):
        p["align"] = ALGN.get(el.get("algn"))
    ln = el.find("a:lnSpc", NS)
    if ln is not None:
        pct, pts = ln.find("a:spcPct", NS), ln.find("a:spcPts", NS)
        if pct is not None:
            p["line_pct"] = int(pct.get("val")) / 100000.0
        elif pts is not None:
            p["line_pts"] = int(pts.get("val")) / 100.0
    for tag, key in (("spcBef", "before_pts"), ("spcAft", "after_pts")):
        sp = el.find("a:" + tag, NS)
        if sp is not None and sp.find("a:spcPts", NS) is not None:
            p[key] = int(sp.find("a:spcPts", NS).get("val")) / 100.0
    if el.find("a:buNone", NS) is not None:
        p["bullet"] = None
    else:
        bc = el.find("a:buChar", NS)
        if bc is not None and bc.get("char"):
            bf = el.find("a:buFont", NS)
            fam = bf.get("typeface") if bf is not None else None
            ch = bc.get("char")[0]
            is_sym = fam and any(s in fam.lower() for s in ("wingding", "webding", "symbol"))
            p["bullet"] = {"char": WINGDINGS.get(ord(ch), ch) if is_sym else ch, "font": None if is_sym else fam}
    df = el.find("a:defRPr", NS)
    if df is not None:
        if df.get("sz"):
            p["size"] = int(df.get("sz")) / 100.0
        if df.get("b") is not None:
            p["bold"] = df.get("b") == "1"
        sr = df.find("a:solidFill/a:srgbClr", NS)
        if sr is not None:
            p["color"] = "#" + sr.get("val").lower()
        lt = df.find("a:latin", NS)
        if lt is not None:
            p["font"] = lt.get("typeface")
    return p


def levels_from(parent):
    out = {}
    if parent is None:
        return out
    for child in parent:
        name = child.tag.split("}")[-1]
        if name.startswith("lvl") and name.endswith("pPr"):
            out[int(name[3]) - 1] = parse_props(child)
    return out


def build_master_levels(master):
    txs = master.element.find(".//p:txStyles", NS)
    out = {"title": {}, "body": {}}
    if txs is not None:
        out["title"] = levels_from(txs.find("p:titleStyle", NS))
        out["body"] = levels_from(txs.find("p:bodyStyle", NS))
    return out


def ph_lststyle_levels(shape):
    try:
        return levels_from(shape.text_frame._txBody.find("a:lstStyle", NS))
    except Exception:
        return {}


def resolve_paragraphs(tf, group, ph_levels):
    mlv = MASTER_LEVELS.get(group, {})
    out = []
    for p in tf.paragraphs:
        lvl = p.level or 0
        eff = {}
        eff.update(mlv.get(lvl, {}))
        eff.update(ph_levels.get(lvl, {}))
        eff.update(parse_props(p._p.find("a:pPr", NS)))
        pdef = {"size": eff.get("size"), "bold": eff.get("bold"),
                "color": eff.get("color"), "font": eff.get("font")}
        runs = []
        for r in p.runs:
            f = r.font
            rp = dict(pdef)
            if f.size is not None:
                rp["size"] = f.size.pt
            if f.bold is not None:
                rp["bold"] = f.bold
            col = rgb_or_none(f.color) if r.text else None
            if col:
                rp["color"] = col
            if f.name:
                rp["font"] = f.name
            runs.append({"t": r.text, "size": rp["size"], "bold": bool(rp["bold"]),
                         "italic": bool(f.italic), "underline": bool(f.underline),
                         "color": rp["color"], "font": rp["font"]})
        out.append({"level": lvl, "align": eff.get("align"),
                    "marL": eff.get("marL", 0.0), "indent": eff.get("indent", 0.0),
                    "line_pct": eff.get("line_pct"), "line_pts": eff.get("line_pts"),
                    "before_pts": eff.get("before_pts"), "after_pts": eff.get("after_pts"),
                    "bullet": eff.get("bullet"), "runs": runs})
    return out


def para_plain(paras):
    return "\n".join("".join(r["t"] for r in p["runs"]) for p in paras).strip()


def dominant_size(paras):
    sizes = [r["size"] for p in paras for r in p["runs"] if r["size"]]
    return max(sizes) if sizes else None


# ----------------------------- image pipeline -----------------------------
def _ensure(d):
    os.makedirs(d, exist_ok=True)
    return d


def save_image(blob, ext, slide_no, idx, assets_dir):
    ext = (ext or "bin").lower()
    base = "s%02d_img%02d" % (slide_no, idx)
    orig_dir = _ensure(os.path.join(assets_dir, "figures", "originals"))
    disp_dir = _ensure(os.path.join(assets_dir, "figures", "display"))
    orig_name = base + "." + ext
    with open(os.path.join(orig_dir, orig_name), "wb") as f:
        f.write(blob)
    nbytes = len(blob)
    rec = {"slide": slide_no, "format": ext, "original_bytes": nbytes,
           "original_path": "assets/figures/originals/" + orig_name,
           "display_path": "assets/figures/originals/" + orig_name,
           "width": None, "height": None, "resized": False, "note": ""}
    needs, big = ext in NEED_RASTER, nbytes > DISPLAY_BIG_BYTES
    if _HAS_PIL and (needs or big):
        try:
            im = Image.open(os.path.join(orig_dir, orig_name)); im.load()
            rec["width"], rec["height"] = im.size
            fmt = "png" if (im.mode in ("RGBA", "P", "LA") or needs) else "jpg"
            out = im
            if max(im.size) > DISPLAY_MAXDIM:
                r = DISPLAY_MAXDIM / max(im.size)
                out = im.resize((int(im.width * r), int(im.height * r)), Image.LANCZOS)
                rec["resized"] = True
            dn = base + "." + fmt
            if fmt == "jpg":
                out.convert("RGB").save(os.path.join(disp_dir, dn), quality=92)
            else:
                out.save(os.path.join(disp_dir, dn))
            rec["display_path"] = "assets/figures/display/" + dn
            rec["note"] = ("rasterized %s->%s" % (ext, fmt)) if needs else "optimized display copy"
        except Exception as e:
            rec["note"] = "display-copy FAILED (%s); original kept" % type(e).__name__
            REPORT["failures"].append({"slide": slide_no, "what": orig_name, "error": str(e)})
    elif _HAS_PIL:
        try:
            with Image.open(os.path.join(orig_dir, orig_name)) as im:
                rec["width"], rec["height"] = im.size
        except Exception:
            pass
    REPORT["images"].append(rec)
    return rec


# ----------------------------- media engine -----------------------------
def handle_video(pic, slide, slide_no, idx, assets_dir):
    vf = pic.find(".//" + A + "videoFile")
    if vf is None:
        return None
    rid = vf.get(R + "link") or vf.get(R + "embed")
    rel = slide.part.rels.get(rid) if rid else None
    media_dir = _ensure(os.path.join(assets_dir, "media", "originals"))
    poster_dir = _ensure(os.path.join(assets_dir, "media", "posters"))
    rec = {"slide": slide_no, "linked": False, "format": None, "original_path": None,
           "display_path": None, "poster_path": None, "playable": False,
           "loop": True, "muted": True, "autoplay": False, "note": ""}
    src_rel = None
    try:
        if rel is None:
            rec["note"] = "video relationship missing"
        elif rel.is_external:
            rec["linked"] = True
            tgt = rel.target_ref
            rec["note"] = "LINKED: " + tgt
            if os.path.exists(tgt):
                ext = tgt.rsplit(".", 1)[-1].lower()
                name = "s%02d_vid%02d.%s" % (slide_no, idx, ext)
                shutil.copyfile(tgt, os.path.join(media_dir, name))
                src_rel, rec["format"] = "assets/media/originals/" + name, ext
            else:
                REPORT["failures"].append({"slide": slide_no, "what": "linked video", "error": "not found: " + tgt})
        else:
            part = rel.target_part
            ext = (part.partname.ext or "bin").lower()
            name = "s%02d_vid%02d.%s" % (slide_no, idx, ext)
            with open(os.path.join(media_dir, name), "wb") as f:
                f.write(part.blob)
            src_rel, rec["format"] = "assets/media/originals/" + name, ext
    except Exception as e:
        rec["note"] = "extraction FAILED: " + type(e).__name__
        REPORT["failures"].append({"slide": slide_no, "what": "video", "error": str(e)})
    rec["original_path"] = src_rel
    if rec["format"] in BROWSER_VIDEO and src_rel:
        rec["playable"], rec["display_path"] = True, src_rel
    elif src_rel:
        rec["note"] = (rec["note"] + " | " if rec["note"] else "") + \
            "format '%s' not browser-playable (no transcoder); poster shown" % rec["format"]
    try:
        blip = pic.find(".//" + A + "blip")
        prid = blip.get(R + "embed") if blip is not None else None
        prel = slide.part.rels.get(prid) if prid else None
        if prel is not None and not prel.is_external:
            pext = (prel.target_part.partname.ext or "png").lower()
            out_ext = "png" if (pext in NEED_RASTER and _HAS_PIL) else pext
            pname = "s%02d_poster%02d.%s" % (slide_no, idx, out_ext)
            ppath = os.path.join(poster_dir, pname)
            if pext in NEED_RASTER and _HAS_PIL:
                im = Image.open(io.BytesIO(prel.target_part.blob)); im.load(); im.save(ppath)
            else:
                with open(ppath, "wb") as f:
                    f.write(prel.target_part.blob)
            rec["poster_path"] = "assets/media/posters/" + pname
    except Exception:
        pass
    REPORT["videos"].append(rec)
    return rec


# ----------------------------- classification -----------------------------
def classify_role(group, geo, text, fsize):
    if group == "title":
        return "title"
    yf = geo["yf"]
    low = yf > 0.9
    if low:
        if any(c.isdigit() for c in text) and (";" in text or "et al" in text.lower()):
            return "citation"
        return "footer"
    return "body" if group == "body" else "annotation"


# ----------------------------- shape walk -----------------------------
def _geo(sh, W, H, xform):
    try:
        l, t, w, h = int(sh.left or 0), int(sh.top or 0), int(sh.width or 0), int(sh.height or 0)
    except Exception:
        l = t = w = h = 0
    if xform:
        l, t, w, h = xform(l, t, w, h)
    return frac(l, t, w, h, W, H)


def walk(shapes, W, H, out, assets_dir, slide, slide_no, xform=None, counters=None):
    counters = counters if counters is not None else {"img": 0, "vid": 0}
    for sh in shapes:
        st = getattr(sh, "shape_type", None)
        if st == MSO_SHAPE_TYPE.GROUP:
            walk(sh.shapes, W, H, out, assets_dir, slide, slide_no, group_child_transform(sh), counters)
            continue
        geo = _geo(sh, W, H, xform)

        if sh._element.tag == P + "pic" and sh._element.find(".//" + A + "videoFile") is not None:
            counters["vid"] += 1
            rec = handle_video(sh._element, slide, slide_no, counters["vid"], assets_dir)
            if rec is not None:
                out.append({"kind": "video", "role": "video", "geo": geo, "media": rec})
            continue

        if st == MSO_SHAPE_TYPE.PICTURE:
            counters["img"] += 1
            try:
                img = sh.image
                rec = save_image(img.blob, img.ext, slide_no, counters["img"], assets_dir)
                src = rec["display_path"]
            except Exception as e:
                src = None
                REPORT["failures"].append({"slide": slide_no, "what": "picture", "error": str(e)})
            small_corner = geo["wf"] < 0.18 and geo["hf"] < 0.18 and geo["yf"] < 0.14
            out.append({"kind": "picture", "role": "logo" if small_corner else "figure",
                        "geo": geo, "src": src})
            continue

        if getattr(sh, "has_table", False):
            out.append({"kind": "table", "role": "table", "geo": geo,
                        "rows": [[c.text for c in row.cells] for row in sh.table.rows]})
            continue

        if st == MSO_SHAPE_TYPE.LINE or sh.__class__.__name__ == "Connector":
            try:
                color = rgb_or_none(sh.line.color)
            except Exception:
                color = None
            out.append({"kind": "line", "role": "bottom-line" if geo["yf"] > 0.85 else "rule",
                        "geo": geo, "color": color})
            continue

        if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip():
            # group: title / body (inherits bodyStyle + bullets) / none (text box)
            group = "none"
            try:
                if sh.is_placeholder:
                    ph = str(sh.placeholder_format.type)
                    if "TITLE" in ph and "SUBTITLE" not in ph:
                        group = "title"
                    elif any(k in ph for k in ("BODY", "OBJECT", "CONTENT", "SUBTITLE")):
                        group = "body"
            except Exception:
                pass
            if geo["yf"] > 0.9:           # bottom band = footer/citation, not bulleted body
                group = "none"
            paras = resolve_paragraphs(sh.text_frame, group, ph_lststyle_levels(sh))
            text, fsize = para_plain(paras), dominant_size(paras)
            role = classify_role(group, geo, text, fsize)
            out.append({"kind": "text", "role": role, "geo": geo, "paras": paras,
                        "plain": text, "size": fsize})
            continue

        if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                fill = rgb_or_none(sh.fill.fore_color)
            except Exception:
                fill = None
            out.append({"kind": "shape",
                        "role": "bottom-line" if (geo["yf"] > 0.85 and geo["hf"] < 0.03) else "generic",
                        "geo": geo, "fill": fill})


# ----------------------------- chrome -----------------------------
def line_color(sh):
    el = sh._element.find(".//a:ln//a:srgbClr", NS)
    return "#" + el.get("val").lower() if el is not None else None


def master_title_style(master):
    lv = MASTER_LEVELS.get("title", {}).get(0, {})
    return {"size_pt": lv.get("size"), "bold": lv.get("bold"), "color": lv.get("color")}


def extract_chrome(prs, W, H, assets_dir):
    chrome = {"logo": None, "bottom_line": None, "footer_area": None, "title_style": None}
    logo_dir = _ensure(os.path.join(assets_dir, "logos"))
    master = prs.slide_masters[0]
    for item in master.shapes:
        st = getattr(item, "shape_type", None)
        if st == MSO_SHAPE_TYPE.PICTURE and chrome["logo"] is None:
            try:
                img = item.image
                name = "logo." + img.ext
                with open(os.path.join(logo_dir, name), "wb") as f:
                    f.write(img.blob)
                chrome["logo"] = {"src": "assets/logos/" + name, "geo": shape_frac(item, W, H)}
                REPORT["logos"].append({"src": "assets/logos/" + name, "native": list(img.size)})
            except Exception as e:
                REPORT["failures"].append({"slide": "master", "what": "logo", "error": str(e)})
        if (st == MSO_SHAPE_TYPE.LINE or item.__class__.__name__ == "Connector") and chrome["bottom_line"] is None:
            g = shape_frac(item, W, H)
            if g["yf"] > 0.8:
                chrome["bottom_line"] = {"geo": g, "color": line_color(item) or "#0070c0"}
    chrome["title_style"] = master_title_style(master)
    return chrome


def infer_design_system(prs, slides, chrome):
    theme_major = theme_minor = None
    try:
        fs = prs.slide_masters[0].element.find(".//a:fontScheme", NS)
        mj, mn = fs.find(".//a:majorFont/a:latin", NS), fs.find(".//a:minorFont/a:latin", NS)
        theme_major = mj.get("typeface") if mj is not None else None
        theme_minor = mn.get("typeface") if mn is not None else None
    except Exception:
        pass
    ts = chrome.get("title_style") or {}
    return {"fonts": {"major": theme_major, "minor": theme_minor},
            "title_size_pt": ts.get("size_pt"),
            "title_color": ts.get("color") if (ts.get("color") and ts.get("color") != "#000000") else "#1f3864",
            "brand_line_color": (chrome.get("bottom_line") or {}).get("color")}


def main():
    if len(sys.argv) < 3:
        print("Usage: python extract_design_system.py <input.pptx> <output_dir>")
        sys.exit(1)
    global MASTER_LEVELS
    src, outdir = sys.argv[1], sys.argv[2]
    assets_dir = _ensure(os.path.join(outdir, "assets"))
    prs = Presentation(src)
    W, H = int(prs.slide_width), int(prs.slide_height)
    aspect = "16:9" if abs(W / H - 16 / 9) < 0.02 else ("4:3" if abs(W / H - 4 / 3) < 0.02 else "%d:%d" % (W, H))
    MASTER_LEVELS = build_master_levels(prs.slide_masters[0])
    chrome = extract_chrome(prs, W, H, assets_dir)

    slides = []
    for i, slide in enumerate(prs.slides):
        objs = []
        walk(slide.shapes, W, H, objs, assets_dir, slide, i + 1)
        slides.append({"number": i + 1, "background": None, "objects": objs})

    extract = {"deck": {"w_emu": W, "h_emu": H, "w_pt": W / EMU_PT, "aspect": aspect, "slides": len(slides)},
               "chrome": chrome, "design_system": infer_design_system(prs, slides, chrome),
               "master_levels": MASTER_LEVELS, "slides": slides}
    with open(os.path.join(outdir, "design-extract.json"), "w", encoding="utf-8") as f:
        json.dump(extract, f, ensure_ascii=False, indent=2)
    with open(os.path.join(assets_dir, "extraction-report.json"), "w", encoding="utf-8") as f:
        json.dump(REPORT, f, ensure_ascii=False, indent=2)

    print("Extracted %d slides. chrome logo=%s line=%s | images=%d videos=%d failures=%d"
          % (len(slides), bool(chrome["logo"]), bool(chrome["bottom_line"]),
             len(REPORT["images"]), len(REPORT["videos"]), len(REPORT["failures"])))
    for v in REPORT["videos"]:
        print("  VIDEO s%d %s playable=%s %s" % (v["slide"], v["format"], v["playable"], v["note"]))


if __name__ == "__main__":
    main()
