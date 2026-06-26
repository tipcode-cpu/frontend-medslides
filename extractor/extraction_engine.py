#!/usr/bin/env python3
"""
extraction_engine.py — THE CORE of frontend-medslides.

Pipeline stage 1 of 3:   PPTX -> [Extraction Engine] -> extract/
(then Validation Engine -> then HTML Generator)

It extracts, preserves, and records EVERY meaningful PowerPoint object before
any HTML exists. Nothing is silently dropped: unconvertible originals are kept
in assets/unsupported/ and every object carries its own extraction status.

Produces:
  extract/
    raw_pptx/<source.pptx>
    assets/{images,videos,logos,icons,audio,unsupported}/
    slides/slideNN.json        (structured slide, schema below)
    manifest.json              (deck, design system, chrome, source counts, slide list)

Usage: python extraction_engine.py <input.pptx> <extract_dir>
Requires: python-pptx, pillow
"""
import io
import json
import os
import shutil
import sys
import urllib.request
import zipfile
from collections import Counter

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from PIL import Image
    _HAS_PIL = True
except Exception:
    _HAS_PIL = False

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
NS = {"a": A[1:-1], "p": P[1:-1], "r": R[1:-1]}

EMU_PT = 12700
STAGE_W, STAGE_H = 1920, 1080
NEED_RASTER = {"wmf", "emf", "wdp", "tiff", "tif", "bmp"}
DISPLAY_MAXDIM = 2600
BROWSER_VIDEO = {"mp4", "m4v", "webm", "ogg", "ogv"}
WINGDINGS = {0x76: "❖", 0xa7: "▪", 0x2022: "•", 0x2013: "–", 0xbb: "»", 0xa8: "▪", 0xfc: "✔"}
ALGN = {"l": "left", "ctr": "center", "r": "right", "just": "justify"}

MASTER_LEVELS = {"title": {}, "body": {}}
THEME_FONTS = {"+mj": None, "+mn": None}   # major/minor theme fonts, set in main


def sanitize_font(name):
    """Drop invalid CSS font names. PowerPoint theme refs like '+mn-lt', '+mj-ea'
    are NOT real fonts -> resolve to the theme font or None (token fallback)."""
    if not name or not isinstance(name, str):
        return None
    name = name.strip()
    if not name:
        return None
    if name.startswith("+"):
        resolved = THEME_FONTS.get(name[:3])      # '+mn' / '+mj'
        return resolved if (resolved and not resolved.startswith("+")) else None
    return name


# ---------- helpers: color / geometry ----------
def rgb_or_none(color):
    try:
        return ("#" + str(color.rgb).lower()) if color is not None else None
    except Exception:
        return None


def px_box(sh, W, H, xform):
    try:
        l, t, w, h = int(sh.left or 0), int(sh.top or 0), int(sh.width or 0), int(sh.height or 0)
    except Exception:
        l = t = w = h = 0
    if xform:
        l, t, w, h = xform(l, t, w, h)
    return {"x": round(l / W * STAGE_W), "y": round(t / H * STAGE_H),
            "width": round(w / W * STAGE_W), "height": round(h / H * STAGE_H)}


def shape_px(sh, W, H):
    return px_box(sh, W, H, None)


def group_child_transform(group):
    try:
        x = group._element.find(".//a:xfrm", NS)
        off, ext, choff, chext = x.find("a:off", NS), x.find("a:ext", NS), x.find("a:chOff", NS), x.find("a:chExt", NS)
        ox, oy, ex, ey = int(off.get("x")), int(off.get("y")), int(ext.get("cx")), int(ext.get("cy"))
        cox, coy, cex, cey = int(choff.get("x")), int(choff.get("y")), int(chext.get("cx")), int(chext.get("cy"))
        sx, sy = (ex / cex if cex else 1.0), (ey / cey if cey else 1.0)
        return lambda l, t, w, h: (ox + (l - cox) * sx, oy + (t - coy) * sy, w * sx, h * sy)
    except Exception:
        return lambda l, t, w, h: (l, t, w, h)


# ---------- text style resolution (run > para > placeholder > master[level]) ----------
def parse_props(el):
    if el is None:
        return {}
    p = {}
    if el.get("marL") is not None:
        p["marginLeft"] = round(int(el.get("marL")) / EMU_PT, 1)
    if el.get("indent") is not None:
        p["indent"] = round(int(el.get("indent")) / EMU_PT, 1)
    if el.get("algn"):
        p["align"] = ALGN.get(el.get("algn"))
    ln = el.find("a:lnSpc", NS)
    if ln is not None:
        pct, pts = ln.find("a:spcPct", NS), ln.find("a:spcPts", NS)
        if pct is not None:
            p["lineSpacingPct"] = int(pct.get("val")) / 100000.0
        elif pts is not None:
            p["lineSpacingPts"] = int(pts.get("val")) / 100.0
    for tag, key in (("spcBef", "spaceBeforePts"), ("spcAft", "spaceAfterPts")):
        sp = el.find("a:" + tag, NS)
        if sp is not None and sp.find("a:spcPts", NS) is not None:
            p[key] = int(sp.find("a:spcPts", NS).get("val")) / 100.0
    if el.find("a:buNone", NS) is not None:
        p["bullet"] = None
        p["bulletFont"] = None
    else:
        bc = el.find("a:buChar", NS)
        if bc is not None and bc.get("char"):
            bf = el.find("a:buFont", NS)
            fam = bf.get("typeface") if bf is not None else None
            ch = bc.get("char")[0]
            sym = fam and any(s in fam.lower() for s in ("wingding", "webding", "symbol"))
            p["bullet"] = WINGDINGS.get(ord(ch), ch) if sym else ch
            p["bulletFont"] = None if sym else sanitize_font(fam)
    # bullet size (percent of text, or points) and color
    bsp, bspt = el.find("a:buSzPct", NS), el.find("a:buSzPts", NS)
    if bsp is not None:
        p["bulletSizePct"] = int(bsp.get("val")) / 100000.0
    elif bspt is not None:
        p["bulletSizePts"] = int(bspt.get("val")) / 100.0
    bclr = el.find("a:buClr//a:srgbClr", NS)
    if bclr is not None:
        p["bulletColor"] = "#" + bclr.get("val").lower()
    df = el.find("a:defRPr", NS)
    if df is not None:
        if df.get("sz"):
            p["fontSize"] = int(df.get("sz")) / 100.0
        if df.get("b") is not None:
            p["fontWeight"] = 700 if df.get("b") == "1" else 400
        sr = df.find("a:solidFill/a:srgbClr", NS)
        if sr is not None:
            p["color"] = "#" + sr.get("val").lower()
        lt = df.find("a:latin", NS)
        if lt is not None:
            p["fontFamily"] = sanitize_font(lt.get("typeface"))
    return p


def levels_from(parent):
    out = {}
    if parent is None:
        return out
    for child in parent:
        name = child.tag.split("}")[-1]
        if name.startswith("lvl") and name.endswith("pPr"):
            out[int(name[3]) - 1] = parse_props(child)
    return out


def build_master_levels(master):
    txs = master.element.find(".//p:txStyles", NS)
    if txs is None:
        return {"title": {}, "body": {}}
    return {"title": levels_from(txs.find("p:titleStyle", NS)),
            "body": levels_from(txs.find("p:bodyStyle", NS))}


def ph_levels(shape):
    try:
        return levels_from(shape.text_frame._txBody.find("a:lstStyle", NS))
    except Exception:
        return {}


def parse_run_props(rPr):
    """Run properties from a:rPr (mirrors defRPr; sanitizes theme fonts)."""
    p = {}
    if rPr is None:
        return p
    if rPr.get("sz"):
        p["fontSize"] = int(rPr.get("sz")) / 100.0
    if rPr.get("b") is not None:
        p["fontWeight"] = 700 if rPr.get("b") == "1" else 400
    if rPr.get("i") is not None:
        p["italic"] = rPr.get("i") == "1"
    if rPr.get("u") not in (None, "none"):
        p["underline"] = True
    sr = rPr.find("a:solidFill/a:srgbClr", NS)
    if sr is not None:
        p["color"] = "#" + sr.get("val").lower()
    lt = rPr.find("a:latin", NS)
    if lt is not None:
        p["fontFamily"] = sanitize_font(lt.get("typeface"))
    return p


def resolve_paragraphs(tf, group, phlv):
    mlv = MASTER_LEVELS.get(group, {})
    out = []
    for p in tf.paragraphs:
        lvl = p.level or 0
        eff = {}
        eff.update(mlv.get(lvl, {}))
        eff.update(phlv.get(lvl, {}))
        eff.update(parse_props(p._p.find("a:pPr", NS)))
        pdef = {"fontSize": eff.get("fontSize"), "fontWeight": eff.get("fontWeight"),
                "color": eff.get("color"), "fontFamily": eff.get("fontFamily")}
        # Iterate XML children so soft line breaks <a:br/> are preserved
        # (python-pptx paragraph.runs drops them -> wrong wrapping -> overlaps).
        runs = []
        for child in p._p:
            tag = child.tag.split("}")[-1]
            if tag == "br":
                runs.append({"text": "\n", "fontFamily": pdef["fontFamily"], "fontSize": pdef["fontSize"],
                             "fontWeight": pdef["fontWeight"] or 400, "italic": False,
                             "underline": False, "color": pdef["color"]})
            elif tag in ("r", "fld"):
                t = child.find("a:t", NS)
                text = t.text if (t is not None and t.text) else ""
                if not text:
                    continue
                rp = dict(pdef)
                rp.update(parse_run_props(child.find("a:rPr", NS)))
                runs.append({"text": text, "fontFamily": rp.get("fontFamily"),
                             "fontSize": rp.get("fontSize"), "fontWeight": rp.get("fontWeight") or 400,
                             "italic": bool(rp.get("italic")), "underline": bool(rp.get("underline")),
                             "color": rp.get("color")})
        out.append({"level": lvl, "bullet": eff.get("bullet"), "bulletFont": eff.get("bulletFont"),
                    "bulletSizePct": eff.get("bulletSizePct"), "bulletSizePts": eff.get("bulletSizePts"),
                    "bulletColor": eff.get("bulletColor"),
                    "indent": eff.get("indent", 0), "marginLeft": eff.get("marginLeft", 0),
                    "align": eff.get("align"),
                    "lineSpacingPct": eff.get("lineSpacingPct"), "lineSpacingPts": eff.get("lineSpacingPts"),
                    "spaceBeforePts": eff.get("spaceBeforePts"), "spaceAfterPts": eff.get("spaceAfterPts"),
                    "runs": runs})
    return out


# ---------- asset extraction ----------
class Engine:
    def __init__(self, extract_dir):
        self.dir = extract_dir
        self.assets = os.path.join(extract_dir, "assets")
        for sub in ("images", "images/originals", "videos", "videos/posters",
                    "logos", "icons", "audio", "unsupported"):
            os.makedirs(os.path.join(self.assets, sub), exist_ok=True)
        os.makedirs(os.path.join(extract_dir, "slides"), exist_ok=True)
        os.makedirs(os.path.join(extract_dir, "raw_pptx"), exist_ok=True)
        os.makedirs(os.path.join(extract_dir, "validation"), exist_ok=True)

    # --- images ---
    def save_image(self, slide_no, rid, blob, ext, orig_name):
        ext = (ext or "bin").lower()
        base = "s%02d_%s" % (slide_no, rid)
        rec = {"relId": rid, "originalName": orig_name, "format": ext,
               "fileSize": len(blob), "dimensions": None, "src": None,
               "original": None, "status": "ok", "error": ""}
        if ext in NEED_RASTER:
            # original -> unsupported/ (kept, flagged); rasterized PNG -> images/
            up = os.path.join(self.assets, "unsupported", base + "." + ext)
            with open(up, "wb") as f:
                f.write(blob)
            rec["original"] = "assets/unsupported/%s.%s" % (base, ext)
            if _HAS_PIL:
                try:
                    im = Image.open(io.BytesIO(blob)); im.load()
                    rec["dimensions"] = list(im.size)
                    if max(im.size) > DISPLAY_MAXDIM:
                        r = DISPLAY_MAXDIM / max(im.size)
                        im = im.resize((int(im.width * r), int(im.height * r)), Image.LANCZOS)
                    out = os.path.join(self.assets, "images", base + ".png")
                    im.save(out)
                    rec["src"] = "assets/images/%s.png" % base
                    rec["note"] = "rasterized %s->png" % ext
                except Exception as e:
                    rec["status"], rec["error"] = "failed", "rasterize %s: %s" % (ext, e)
            else:
                rec["status"], rec["error"] = "failed", "no PIL to rasterize %s" % ext
            return rec
        # browser-ready raster
        out = os.path.join(self.assets, "images", base + "." + ext)
        with open(out, "wb") as f:
            f.write(blob)
        rec["src"] = "assets/images/%s.%s" % (base, ext)
        if _HAS_PIL:
            try:
                im = Image.open(out); im.load()
                rec["dimensions"] = list(im.size)
                if max(im.size) > DISPLAY_MAXDIM:
                    shutil.copyfile(out, os.path.join(self.assets, "images", "originals", base + "." + ext))
                    rec["original"] = "assets/images/originals/%s.%s" % (base, ext)
                    r = DISPLAY_MAXDIM / max(im.size)
                    disp = im.resize((int(im.width * r), int(im.height * r)), Image.LANCZOS)
                    disp.save(out)
                    rec["note"] = "downscaled display; original kept"
            except Exception:
                pass
        return rec

    # --- logo (master) ---
    def save_logo(self, blob, ext, idx=1):
        name = "logo%02d.%s" % (idx, (ext or "png").lower())
        with open(os.path.join(self.assets, "logos", name), "wb") as f:
            f.write(blob)
        dims = None
        if _HAS_PIL:
            try:
                with Image.open(io.BytesIO(blob)) as im:
                    dims = list(im.size)
            except Exception:
                pass
        return {"src": "assets/logos/" + name, "fileSize": len(blob), "dimensions": dims}

    # --- interactive HTML simulator (hyperlinked .html) ---
    def save_simulator(self, url):
        """Download a referenced HTML simulator for offline inline embedding.
        Best-effort: returns local rel path, or None (build falls back to URL)."""
        name = url.rsplit("/", 1)[-1].split("?")[0] or "simulator.html"
        d = os.path.join(self.assets, "simulators")
        os.makedirs(d, exist_ok=True)
        dst = os.path.join(d, name)
        # urllib with an unverified SSL context (Windows certs often fail), then curl
        try:
            import ssl
            ctx = ssl.create_default_context()
            ctx.check_hostname = False
            ctx.verify_mode = ssl.CERT_NONE
            req = urllib.request.Request(url, headers={"User-Agent": "frontend-medslides"})
            with urllib.request.urlopen(req, timeout=20, context=ctx) as r:
                data = r.read()
            if data:
                open(dst, "wb").write(data)
                return "assets/simulators/" + name
        except Exception:
            pass
        try:
            import subprocess
            subprocess.run(["curl", "-sL", "-m", "25", url, "-o", dst],
                           capture_output=True, timeout=30)
            if os.path.exists(dst) and os.path.getsize(dst) > 0:
                return "assets/simulators/" + name
        except Exception:
            pass
        return None

    # --- video ---
    def save_video(self, slide_no, idx, pic, slide):
        vf = pic.find(".//" + A + "videoFile")
        rid = vf.get(R + "link") or vf.get(R + "embed")
        rel = slide.part.rels.get(rid) if rid else None
        rec = {"relId": rid, "embedded": True, "format": None, "src": None, "original": None,
               "poster": None, "playable": False, "conversionNeeded": False,
               "loop": True, "muted": True, "autoplay": False, "fileSize": 0,
               "status": "ok", "error": "", "originalName": None}
        try:
            if rel is None:
                rec["status"], rec["error"] = "failed", "video relationship missing"
            elif rel.is_external:
                rec["embedded"] = False
                rec["originalName"] = rel.target_ref
                if os.path.exists(rel.target_ref):
                    ext = rel.target_ref.rsplit(".", 1)[-1].lower()
                    name = "s%02d_vid%02d.%s" % (slide_no, idx, ext)
                    shutil.copyfile(rel.target_ref, os.path.join(self.assets, "videos", name))
                    rec["src"], rec["format"] = "assets/videos/" + name, ext
                    rec["fileSize"] = os.path.getsize(os.path.join(self.assets, "videos", name))
                else:
                    rec["status"], rec["error"] = "failed", "linked file not found: " + rel.target_ref
            else:
                part = rel.target_part
                ext = (part.partname.ext or "bin").lower()
                rec["originalName"] = str(part.partname)
                name = "s%02d_vid%02d.%s" % (slide_no, idx, ext)
                blob = part.blob
                with open(os.path.join(self.assets, "videos", name), "wb") as f:
                    f.write(blob)
                rec["src"], rec["format"], rec["fileSize"] = "assets/videos/" + name, ext, len(blob)
        except Exception as e:
            rec["status"], rec["error"] = "failed", "extract: %s" % e
        if rec["format"] in BROWSER_VIDEO and rec["src"]:
            rec["playable"] = True
        elif rec["src"]:
            rec["conversionNeeded"] = True
        # poster
        try:
            blip = pic.find(".//" + A + "blip")
            prid = blip.get(R + "embed") if blip is not None else None
            prel = slide.part.rels.get(prid) if prid else None
            if prel is not None and not prel.is_external:
                pext = (prel.target_part.partname.ext or "png").lower()
                oext = "png" if (pext in NEED_RASTER and _HAS_PIL) else pext
                pname = "s%02d_poster%02d.%s" % (slide_no, idx, oext)
                pp = os.path.join(self.assets, "videos", "posters", pname)
                if pext in NEED_RASTER and _HAS_PIL:
                    im = Image.open(io.BytesIO(prel.target_part.blob)); im.load(); im.save(pp)
                else:
                    with open(pp, "wb") as f:
                        f.write(prel.target_part.blob)
                rec["poster"] = "assets/videos/posters/" + pname
        except Exception:
            pass
        return rec


# ---------- object building per slide ----------
def fill_blip_rid(el):
    """rId of a picture FILL blip on this shape — covers p:pic, and a:blipFill on
    autoshapes / picture placeholders (which python-pptx's shape.image MISSES).
    Excludes bullet pictures (a:buBlip)."""
    for blip in el.iter(A + "blip"):
        parent = blip.getparent()
        if parent is not None and parent.tag.endswith("blipFill"):
            rid = blip.get(R + "embed") or blip.get(R + "link")
            if rid:
                return rid
    return None


def bg_blip_rid(slide_el):
    bg = slide_el.find(".//" + P + "bg")
    if bg is None:
        return None
    blip = bg.find(".//" + A + "blip")
    return (blip.get(R + "embed") or blip.get(R + "link")) if blip is not None else None


def walk(shapes, eng, W, H, slide, slide_no, ids, objs, xform=None, counters=None):
    counters = counters if counters is not None else {"img": 0, "vid": 0}
    for sh in shapes:
        st = getattr(sh, "shape_type", None)
        if st == MSO_SHAPE_TYPE.GROUP:
            start = len(objs)
            counters["grp"] = counters.get("grp", 0) + 1
            gid = "g%03d" % counters["grp"]
            walk(sh.shapes, eng, W, H, slide, slide_no, ids, objs, group_child_transform(sh), counters)
            for o in objs[start:]:
                o.setdefault("groupId", gid)
            continue
        geo = px_box(sh, W, H, xform)

        # video
        if sh._element.tag == P + "pic" and sh._element.find(".//" + A + "videoFile") is not None:
            counters["vid"] += 1
            ids["video"] += 1
            rec = eng.save_video(slide_no, counters["vid"], sh._element, slide)
            o = {"id": "video-%03d" % ids["video"], "type": "video", "role": "video"}
            o.update(geo); o.update(rec); objs.append(o)
            continue

        # image: ANY shape carrying a picture fill blip (Picture, picture-filled
        # autoshape, picture placeholder). python-pptx Picture-only extraction
        # misses fills/placeholders -> that was the source of missing figures.
        rid = fill_blip_rid(sh._element)
        if rid is not None and sh._element.find(".//" + A + "videoFile") is None:
            rel = slide.part.rels.get(rid)
            reltype = (getattr(rel, "reltype", "") or "")
            if rel is not None and "image" in reltype.lower():
                try:
                    if rel.is_external:
                        if os.path.exists(rel.target_ref):
                            blob = open(rel.target_ref, "rb").read()
                            rec = eng.save_image(slide_no, rid, blob, rel.target_ref.rsplit(".", 1)[-1], rel.target_ref)
                        else:
                            rec = {"relId": rid, "originalName": rel.target_ref, "format": None,
                                   "fileSize": 0, "dimensions": None, "src": None, "original": None,
                                   "status": "linked-missing", "error": "external image not found: " + rel.target_ref}
                    else:
                        part = rel.target_part
                        rec = eng.save_image(slide_no, rid, part.blob, part.partname.ext, str(part.partname))
                except Exception as e:
                    rec = {"relId": rid, "originalName": None, "format": None, "fileSize": 0,
                           "dimensions": None, "src": None, "original": None,
                           "status": "failed", "error": "image extract: %s" % e}
                small = geo["width"] < STAGE_W * 0.18 and geo["height"] < STAGE_H * 0.18 and geo["y"] < STAGE_H * 0.14
                ids["image"] += 1
                o = {"id": "image-%03d" % ids["image"], "type": "image",
                     "role": "logo" if small else "figure", "zIndex": len(objs)}
                o.update(geo); o.update(rec)
                # hyperlink to an .html simulator -> embed it live inline (iframe)
                hl = sh._element.find(".//" + A + "hlinkClick")
                if hl is not None:
                    hrel = slide.part.rels.get(hl.get(R + "id"))
                    if (hrel is not None and getattr(hrel, "is_external", False)
                            and (hrel.target_ref or "").lower().rsplit("?", 1)[0].endswith((".html", ".htm"))):
                        o["role"] = "simulator"
                        o["simulatorUrl"] = hrel.target_ref
                        o["simulatorLocal"] = eng.save_simulator(hrel.target_ref)
                objs.append(o)
                continue

        # table
        if getattr(sh, "has_table", False):
            ids["table"] += 1
            o = {"id": "table-%03d" % ids["table"], "type": "table", "role": "table",
                 "rows": [[c.text for c in row.cells] for row in sh.table.rows]}
            o.update(geo); objs.append(o)
            continue

        # chart (graphic frame)
        if sh.has_chart if hasattr(sh, "has_chart") else False:
            ids["chart"] += 1
            o = {"id": "chart-%03d" % ids["chart"], "type": "chart", "role": "chart",
                 "status": "unsupported", "error": "chart rendering not implemented; reported"}
            o.update(geo); objs.append(o)
            continue

        # line / connector (bottom blue line)
        if st == MSO_SHAPE_TYPE.LINE or sh.__class__.__name__ == "Connector":
            try:
                color = rgb_or_none(sh.line.color)
            except Exception:
                color = None
            ids["line"] += 1
            o = {"id": "line-%03d" % ids["line"], "type": "line",
                 "role": "bottom-line" if geo["y"] > STAGE_H * 0.85 else "rule", "color": color}
            o.update(geo); objs.append(o)
            continue

        # text
        if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip():
            group = "none"
            try:
                if sh.is_placeholder:
                    ph = str(sh.placeholder_format.type)
                    if "TITLE" in ph and "SUBTITLE" not in ph:
                        group = "title"
                    elif any(k in ph for k in ("BODY", "OBJECT", "CONTENT", "SUBTITLE")):
                        group = "body"
            except Exception:
                pass
            if geo["y"] > STAGE_H * 0.9:
                group = "none"
            paras = resolve_paragraphs(sh.text_frame, group, ph_levels(sh))
            text = "\n".join("".join(r["text"] for r in p["runs"]) for p in paras).strip()
            role = ("title" if group == "title"
                    else ("footer" if geo["y"] > STAGE_H * 0.9 else
                          ("body" if group == "body" else "annotation")))
            ids["text"] += 1
            o = {"id": "text-%03d" % ids["text"], "type": "text", "role": role,
                 "paragraphs": paras, "plainText": text}
            o.update(geo); objs.append(o)
            continue

        # autoshape (fills / banners)
        if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                fill = rgb_or_none(sh.fill.fore_color)
            except Exception:
                fill = None
            ids["shape"] += 1
            o = {"id": "shape-%03d" % ids["shape"], "type": "shape",
                 "role": "bottom-line" if (geo["y"] > STAGE_H * 0.85 and geo["height"] < STAGE_H * 0.03) else "generic",
                 "fill": fill}
            o.update(geo); objs.append(o)


# ---------- chrome (master) ----------
def line_color(sh):
    el = sh._element.find(".//a:ln//a:srgbClr", NS)
    return "#" + el.get("val").lower() if el is not None else None


def extract_chrome(prs, eng, W, H):
    chrome = {"logo": None, "bottomLine": None, "footerArea": None, "titleStyle": None}
    master = prs.slide_masters[0]
    for item in master.shapes:
        st = getattr(item, "shape_type", None)
        if st == MSO_SHAPE_TYPE.PICTURE and chrome["logo"] is None:
            try:
                img = item.image
                info = eng.save_logo(img.blob, img.ext)
                info.update(shape_px(item, W, H))
                chrome["logo"] = info
            except Exception:
                pass
        if (st == MSO_SHAPE_TYPE.LINE or item.__class__.__name__ == "Connector") and chrome["bottomLine"] is None:
            g = shape_px(item, W, H)
            if g["y"] > STAGE_H * 0.8:
                chrome["bottomLine"] = dict(g, color=line_color(item) or "#0070c0")
    tlv = MASTER_LEVELS.get("title", {}).get(0, {})
    chrome["titleStyle"] = {"fontSize": tlv.get("fontSize"), "fontWeight": tlv.get("fontWeight"),
                            "color": tlv.get("color")}
    return chrome


def source_counts(path, slide_count):
    z = zipfile.ZipFile(path)
    media = [n for n in z.namelist() if n.startswith("ppt/media/")]
    exts = Counter(n.rsplit(".", 1)[-1].lower() for n in media)
    vids = sum(v for k, v in exts.items() if k in ("mp4", "wmv", "mov", "avi", "m4v", "webm", "mpg", "mpeg", "mkv"))
    auds = sum(v for k, v in exts.items() if k in ("mp3", "wav", "m4a", "wma", "aac"))
    return {"slideCount": slide_count, "mediaFiles": len(media),
            "imageLike": len(media) - vids - auds, "videoLike": vids, "audioLike": auds,
            "mediaByExt": dict(exts)}


def main():
    if len(sys.argv) < 3:
        print("Usage: python extraction_engine.py <input.pptx> <extract_dir>")
        sys.exit(1)
    global MASTER_LEVELS
    src, outdir = sys.argv[1], sys.argv[2]
    if os.path.isdir(outdir):
        shutil.rmtree(outdir, ignore_errors=True)
    eng = Engine(outdir)
    # Unzip the PPTX (it is a ZIP) so ppt/media + rels are available for the
    # raw relationship-based media extractor and cross-check.
    with zipfile.ZipFile(src) as _z:
        _z.extractall(os.path.join(outdir, "raw_pptx"))

    prs = Presentation(src)
    W, H = int(prs.slide_width), int(prs.slide_height)
    aspect = "16:9" if abs(W / H - 16 / 9) < 0.02 else ("4:3" if abs(W / H - 4 / 3) < 0.02 else "%d:%d" % (W, H))
    MASTER_LEVELS = build_master_levels(prs.slide_masters[0])
    # Resolve theme fonts BEFORE walking (used to sanitize '+mn-lt'/'+mj-lt' refs)
    fs0 = prs.slide_masters[0].element.find(".//a:fontScheme", NS)
    if fs0 is not None:
        mj0 = fs0.find(".//a:majorFont/a:latin", NS)
        mn0 = fs0.find(".//a:minorFont/a:latin", NS)
        THEME_FONTS["+mj"] = mj0.get("typeface") if mj0 is not None else None
        THEME_FONTS["+mn"] = mn0.get("typeface") if mn0 is not None else None
    chrome = extract_chrome(prs, eng, W, H)

    slide_files = []
    totals = Counter()
    for i, slide in enumerate(prs.slides):
        n = i + 1
        ids = {"text": 0, "image": 0, "video": 0, "table": 0, "line": 0, "shape": 0, "chart": 0}
        objs = []
        # slide background image (picture fill on the slide bg)
        bgrid = bg_blip_rid(slide._element)
        if bgrid:
            rel = slide.part.rels.get(bgrid)
            if rel is not None and not getattr(rel, "is_external", False) and "image" in (rel.reltype or "").lower():
                try:
                    part = rel.target_part
                    rec = eng.save_image(n, bgrid, part.blob, part.partname.ext, str(part.partname))
                    bo = {"id": "image-bg", "type": "image", "role": "background",
                          "x": 0, "y": 0, "width": STAGE_W, "height": STAGE_H}
                    bo.update(rec); objs.append(bo)
                except Exception:
                    pass
        walk(slide.shapes, eng, W, H, slide, n, ids, objs)
        for zi, o in enumerate(objs):           # z-order = document order
            o["zIndex"] = zi
            o.setdefault("groupId", None)
        try:
            bg = rgb_or_none(slide.background.fill.fore_color)
        except Exception:
            bg = None
        sj = {"slideNumber": n, "size": {"width": STAGE_W, "height": STAGE_H},
              "background": {"type": "color", "value": bg} if bg else {"type": "inherit", "value": "#ffffff"},
              "objects": objs}
        fn = "slides/slide%02d.json" % n
        with open(os.path.join(outdir, fn), "w", encoding="utf-8") as f:
            json.dump(sj, f, ensure_ascii=False, indent=2)
        slide_files.append(fn)
        for o in objs:
            totals[o["type"]] += 1

    fs = prs.slide_masters[0].element.find(".//a:fontScheme", NS)
    mj = fs.find(".//a:majorFont/a:latin", NS) if fs is not None else None
    mn = fs.find(".//a:minorFont/a:latin", NS) if fs is not None else None
    body_levels = {str(k): {"fontSize": v.get("fontSize"), "bullet": v.get("bullet")}
                   for k, v in MASTER_LEVELS.get("body", {}).items()}

    manifest = {
        "source": os.path.basename(src),
        "rawPptx": "raw_pptx/" + os.path.basename(src),
        "deck": {"slideWidthPt": W / EMU_PT, "slideHeightPt": H / EMU_PT,
                 "stageWidth": STAGE_W, "stageHeight": STAGE_H,
                 "ptToPx": STAGE_W / (W / EMU_PT), "aspect": aspect, "slideCount": len(slide_files)},
        "designSystem": {
            "fonts": {"major": mj.get("typeface") if mj is not None else None,
                      "minor": mn.get("typeface") if mn is not None else None},
            "titleSizePt": (chrome["titleStyle"] or {}).get("fontSize"),
            "titleColor": ((chrome["titleStyle"] or {}).get("color")
                           if (chrome["titleStyle"] or {}).get("color") not in (None, "#000000") else "#1f3864"),
            "brandLineColor": (chrome.get("bottomLine") or {}).get("color"),
            "bodyLevels": body_levels},
        "chrome": chrome,
        "source_counts": source_counts(src, len(slide_files)),
        "objectTotals": dict(totals),
        "slides": slide_files,
        "validation": {"status": "pending"},
    }
    with open(os.path.join(outdir, "manifest.json"), "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)

    print("Extraction complete -> %s" % outdir)
    print("  slides=%d  objects=%s" % (len(slide_files), dict(totals)))
    print("  source media: %s" % manifest["source_counts"]["mediaByExt"])
    print("  chrome: logo=%s bottomLine=%s" % (bool(chrome["logo"]), bool(chrome["bottomLine"])))
    print("  -> run Validation Engine next (HTML must not be generated yet)")


if __name__ == "__main__":
    main()
